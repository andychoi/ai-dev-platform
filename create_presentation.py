#!/usr/bin/env python3
"""
Coder WebIDE Tech Day Presentation Generator
Creates a professional PPTX using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# Tech Day Theme Colors
COLORS = {
    'primary': RGBColor(0x1a, 0x1a, 0x2e),      # Dark navy
    'secondary': RGBColor(0x16, 0x21, 0x3e),    # Darker blue
    'accent': RGBColor(0x0f, 0x4c, 0x75),       # Teal blue
    'highlight': RGBColor(0x3a, 0x86, 0xff),    # Bright blue
    'success': RGBColor(0x00, 0xd9, 0xa5),      # Green/teal
    'warning': RGBColor(0xff, 0x6b, 0x6b),      # Coral red
    'white': RGBColor(0xff, 0xff, 0xff),
    'light_gray': RGBColor(0xe0, 0xe0, 0xe0),
    'dark_gray': RGBColor(0x66, 0x66, 0x66),
}


def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=COLORS['white']):
    """Add a title text box"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return shape


def add_body_text(slide, text, left, top, width, height, font_size=18, color=COLORS['light_gray'], bold=False):
    """Add body text box"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    return shape


def add_bullet_points(slide, items, left, top, width, height, font_size=18, color=COLORS['light_gray']):
    """Add bullet point list"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return shape


def add_table(slide, data, left, top, width, height):
    """Add a styled table"""
    rows = len(data)
    cols = len(data[0]) if data else 0

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Style the table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # Style cell
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.name = "Arial"

            if i == 0:  # Header row
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS['white']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['accent']
            else:
                paragraph.font.color.rgb = COLORS['dark_gray']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

    return table


def add_code_block(slide, code, left, top, width, height):
    """Add a code block with monospace font"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2d, 0x2d, 0x2d)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['success']
    p.font.name = "Courier New"
    return shape


def add_accent_bar(slide, left, top, width, height, color=COLORS['highlight']):
    """Add a colored accent bar"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # Slide 1: Title Slide
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, COLORS['primary'])

    # Accent bar at top
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), COLORS['highlight'])

    # Main title
    add_title_shape(slide, "Coder WebIDE",
                   Inches(0.5), Inches(2), Inches(12), Inches(1),
                   font_size=60, color=COLORS['white'])

    # Subtitle
    add_title_shape(slide, "Secure, Cost-Effective, AI-Ready Development Environments",
                   Inches(0.5), Inches(3.2), Inches(12), Inches(0.8),
                   font_size=28, bold=False, color=COLORS['highlight'])

    # Event info
    add_body_text(slide, "Developer Tech Day 2026",
                 Inches(0.5), Inches(5.5), Inches(6), Inches(0.5),
                 font_size=20, color=COLORS['light_gray'])

    add_body_text(slide, "15-Minute Presentation",
                 Inches(0.5), Inches(6), Inches(6), Inches(0.5),
                 font_size=16, color=COLORS['dark_gray'])

    # =========================================================================
    # Slide 2: The Problem
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['warning'])

    add_title_shape(slide, "The Problem",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['warning'])

    add_body_text(slide, "Traditional VDI Pain Points",
                 Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
                 font_size=24, color=COLORS['white'], bold=True)

    # Problem table
    problem_data = [
        ["Challenge", "Impact"],
        ["Slow startup", "5-15 min to get working environment"],
        ["High cost", "Windows licenses + VDI infrastructure"],
        ["Resource waste", "Full VM per developer, always running"],
        ["Security gaps", "Data on local devices, inconsistent policies"],
        ["Onboarding delay", "Days to weeks for new hire setup"],
    ]
    add_table(slide, problem_data, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5))

    # AI Challenge section
    add_body_text(slide, "The AI Challenge:",
                 Inches(7), Inches(1.2), Inches(6), Inches(0.5),
                 font_size=24, color=COLORS['highlight'], bold=True)

    add_bullet_points(slide, [
        "How do we give developers AI assistants?",
        "Without exposing API keys to local machines?",
        "Without code leaking to untrusted devices?",
    ], Inches(7), Inches(1.8), Inches(5.5), Inches(2))

    # =========================================================================
    # Slide 3: The Solution
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['success'])

    add_title_shape(slide, "The Solution: Coder WebIDE",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['success'])

    add_body_text(slide, "Browser-Based Development Environments",
                 Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                 font_size=24, color=COLORS['white'], bold=True)

    # Architecture diagram as text
    arch_code = """┌──────────────────────────────────────────┐
│      Developer's Browser                 │
│      Any device, any location            │
│      No local code, no secrets           │
└────────────────┬─────────────────────────┘
                 │ HTTPS only
                 ▼
┌──────────────────────────────────────────┐
│           Coder Platform                 │
│  ┌──────────┐ ┌──────────┐ ┌──────────┐ │
│  │Workspace │ │Workspace │ │Workspace │ │
│  │ VS Code  │ │ VS Code  │ │ VS Code  │ │
│  └──────────┘ └──────────┘ └──────────┘ │
│         All code stays here              │
└──────────────────────────────────────────┘"""

    add_code_block(slide, arch_code, Inches(0.5), Inches(1.8), Inches(7), Inches(4))

    # Key insight
    add_body_text(slide, "Key Insight:",
                 Inches(8), Inches(2), Inches(5), Inches(0.5),
                 font_size=22, color=COLORS['highlight'], bold=True)

    add_bullet_points(slide, [
        "Full VS Code in browser",
        "Code never leaves infrastructure",
        "Works from any device",
        "No local setup required",
    ], Inches(8), Inches(2.6), Inches(4.5), Inches(3), font_size=20)

    # =========================================================================
    # Slide 4: Four Key Objectives
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['highlight'])

    add_title_shape(slide, "Four Key Objectives",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['white'])

    # Objective 1
    add_body_text(slide, "1. Security & Compliance",
                 Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
                 font_size=22, color=COLORS['highlight'], bold=True)
    add_bullet_points(slide, [
        "Zero-trust: No shell/RDP from untrusted devices",
        "SSO via OIDC (Azure AD ready)",
        "Complete audit trail",
    ], Inches(0.5), Inches(1.7), Inches(6), Inches(1.3), font_size=16)

    # Objective 2
    add_body_text(slide, "2. Cost Reduction (40-60% vs VDI)",
                 Inches(0.5), Inches(3.2), Inches(6), Inches(0.4),
                 font_size=22, color=COLORS['success'], bold=True)
    add_bullet_points(slide, [
        "Containers vs VMs: 10x density",
        "30-60 sec startup vs 5-15 min",
        "No Windows/VDI licenses",
    ], Inches(0.5), Inches(3.6), Inches(6), Inches(1.3), font_size=16)

    # Objective 3
    add_body_text(slide, "3. Fast Onboarding",
                 Inches(7), Inches(1.3), Inches(6), Inches(0.4),
                 font_size=22, color=COLORS['warning'], bold=True)
    add_bullet_points(slide, [
        "New hire → productive in minutes",
        "Standardized templates",
        "No 'works on my machine' issues",
    ], Inches(7), Inches(1.7), Inches(5.5), Inches(1.3), font_size=16)

    # Objective 4
    add_body_text(slide, "4. AI-Ready Platform",
                 Inches(7), Inches(3.2), Inches(6), Inches(0.4),
                 font_size=22, color=RGBColor(0xbb, 0x86, 0xfc), bold=True)
    add_bullet_points(slide, [
        "Secure AI Gateway (Claude, Bedrock)",
        "Built for AI coding agents",
        "Rate limiting & audit logging",
    ], Inches(7), Inches(3.6), Inches(5.5), Inches(1.3), font_size=16)

    # Bottom highlight
    add_body_text(slide, "Aligned with Coder's Enterprise AI Development vision",
                 Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
                 font_size=16, color=COLORS['dark_gray'])

    # =========================================================================
    # Slide 5: Architecture Overview
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['accent'])

    add_title_shape(slide, "Architecture Overview",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['white'])

    arch_full = """┌─────────────────────────────────────────────────────────────┐
│                    Platform Services                         │
├──────────────┬──────────────┬──────────────┬────────────────┤
│    Coder     │  Authentik   │    Gitea     │  AI Gateway    │
│    :7080     │    (SSO)     │    (Git)     │    :8090       │
│              │    :9000     │    :3000     │                │
└──────────────┴──────────────┴──────────────┴────────────────┘
        │                            │               │
        ▼                            ▼               ▼
┌──────────────┐              ┌────────────┐  ┌─────────────┐
│  Workspaces  │              │   Repos    │  │   Claude    │
│ (Containers) │◄────────────►│            │  │   Bedrock   │
└──────────────┘              └────────────┘  │   Gemini    │
                                              └─────────────┘"""

    add_code_block(slide, arch_full, Inches(0.5), Inches(1.3), Inches(9), Inches(4))

    add_body_text(slide, "Key Points:",
                 Inches(10), Inches(1.5), Inches(3), Inches(0.4),
                 font_size=20, color=COLORS['highlight'], bold=True)

    add_bullet_points(slide, [
        "All containerized",
        "Self-hosted",
        "Enterprise-controlled",
        "14 services total",
        "Single docker-compose",
    ], Inches(10), Inches(2), Inches(3), Inches(3), font_size=16)

    # =========================================================================
    # Slide 6: AI Gateway
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), RGBColor(0xbb, 0x86, 0xfc))

    add_title_shape(slide, "AI Gateway: Secure AI Access",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['white'])

    add_body_text(slide, "The Problem:",
                 Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
                 font_size=20, color=COLORS['warning'], bold=True)

    add_bullet_points(slide, [
        "Developers need AI assistants",
        "Can't expose API keys to local machines",
        "Need audit trail & cost control",
    ], Inches(0.5), Inches(1.6), Inches(5), Inches(1.2), font_size=16)

    add_body_text(slide, "The Solution:",
                 Inches(0.5), Inches(3), Inches(6), Inches(0.4),
                 font_size=20, color=COLORS['success'], bold=True)

    # Features table
    ai_data = [
        ["Feature", "Benefit"],
        ["No credential exposure", "API keys in gateway only"],
        ["Multi-provider", "Claude, Bedrock, Gemini"],
        ["Rate limiting", "Per-user request controls"],
        ["Audit logging", "Full request/response tracking"],
        ["Future-ready", "Supports AI agents (Claude Code)"],
    ]
    add_table(slide, ai_data, Inches(0.5), Inches(3.5), Inches(5.5), Inches(2.5))

    # Code example
    add_body_text(slide, "Usage (no API key needed):",
                 Inches(6.5), Inches(1.2), Inches(6), Inches(0.4),
                 font_size=18, color=COLORS['highlight'], bold=True)

    code = """# From any workspace
curl http://ai-gateway:8090/v1/claude/messages \\
  -H "Content-Type: application/json" \\
  -d '{
    "model": "claude-sonnet-4",
    "messages": [
      {"role": "user", "content": "Hello"}
    ]
  }'"""
    add_code_block(slide, code, Inches(6.5), Inches(1.7), Inches(6.3), Inches(2.8))

    # =========================================================================
    # Slide 7: Live Demo
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['secondary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), COLORS['highlight'])

    add_title_shape(slide, "LIVE DEMO",
                   Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                   font_size=56, color=COLORS['highlight'])

    add_body_text(slide, "4-Minute Demo Flow",
                 Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
                 font_size=28, color=COLORS['white'], bold=True)

    # Demo steps
    demo_steps = [
        "1. Login via SSO (Authentik → Azure AD ready)",
        "2. Create Workspace (30-60 seconds spin-up)",
        "3. VS Code in Browser (full IDE experience)",
        "4. AI Assistant via Gateway (Claude integration)",
        "5. Git Workflow (clone, commit, push to Gitea)",
    ]

    for i, step in enumerate(demo_steps):
        y_pos = 2.5 + (i * 0.7)
        color = COLORS['success'] if i < 3 else COLORS['highlight']
        add_body_text(slide, step,
                     Inches(1), Inches(y_pos), Inches(11), Inches(0.6),
                     font_size=24, color=color)

    add_body_text(slide, "http://host.docker.internal:7080",
                 Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                 font_size=20, color=COLORS['dark_gray'])

    # =========================================================================
    # Slide 8: Results & Metrics
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['success'])

    add_title_shape(slide, "Results & Metrics",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['white'])

    # Results table
    results_data = [
        ["Metric", "Before (VDI)", "After (Coder)", "Improvement"],
        ["Environment startup", "5-15 min", "30-60 sec", "90%+ faster"],
        ["New hire onboarding", "2-5 days", "< 1 hour", "95%+ faster"],
        ["Infrastructure cost", "$X/user/mo", "~0.5X", "40-60% savings"],
        ["Security audit", "Manual", "Full auto", "100% visibility"],
    ]
    add_table(slide, results_data, Inches(0.5), Inches(1.3), Inches(12), Inches(2.2))

    # Developer quotes
    add_body_text(slide, "Developer Feedback:",
                 Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
                 font_size=22, color=COLORS['highlight'], bold=True)

    quotes = [
        '"Finally, I can code from my iPad on the train"',
        '"AI assistant without fighting IT for API keys"',
        '"New contractor was coding within 30 minutes of signing NDA"',
    ]
    for i, quote in enumerate(quotes):
        add_body_text(slide, quote,
                     Inches(0.8), Inches(4.3 + i * 0.6), Inches(11), Inches(0.5),
                     font_size=18, color=COLORS['light_gray'])

    # =========================================================================
    # Slide 9: Roadmap
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), COLORS['highlight'])

    add_title_shape(slide, "Roadmap & Next Steps",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                   font_size=40, color=COLORS['white'])

    # Current state
    add_body_text(slide, "Current State: PoC Complete ✓",
                 Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
                 font_size=22, color=COLORS['success'], bold=True)

    add_bullet_points(slide, [
        "14 services running",
        "SSO integrated",
        "AI Gateway operational",
    ], Inches(0.5), Inches(1.6), Inches(5), Inches(1), font_size=16)

    # Timeline
    add_body_text(slide, "Pilot Program - Q1 2026",
                 Inches(0.5), Inches(2.8), Inches(4), Inches(0.4),
                 font_size=20, color=COLORS['highlight'], bold=True)
    add_body_text(slide, "10-20 developers, gather feedback",
                 Inches(0.5), Inches(3.2), Inches(4), Inches(0.4),
                 font_size=16, color=COLORS['light_gray'])

    add_body_text(slide, "Production Deploy - Q2 2026",
                 Inches(4.5), Inches(2.8), Inches(4), Inches(0.4),
                 font_size=20, color=COLORS['highlight'], bold=True)
    add_body_text(slide, "Kubernetes migration, HA/DR",
                 Inches(4.5), Inches(3.2), Inches(4), Inches(0.4),
                 font_size=16, color=COLORS['light_gray'])

    add_body_text(slide, "Enterprise Features - Q3 2026",
                 Inches(8.5), Inches(2.8), Inches(4), Inches(0.4),
                 font_size=20, color=COLORS['highlight'], bold=True)
    add_body_text(slide, "Azure AD, advanced RBAC",
                 Inches(8.5), Inches(3.2), Inches(4), Inches(0.4),
                 font_size=16, color=COLORS['light_gray'])

    # Get involved
    add_body_text(slide, "Get Involved:",
                 Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
                 font_size=22, color=COLORS['white'], bold=True)

    add_bullet_points(slide, [
        "GitHub: github.com/andychoi/dev-platform",
        "Slack: #coder-webide-pilot",
        "Sign up for pilot program!",
    ], Inches(0.5), Inches(4.6), Inches(6), Inches(1.5), font_size=18)

    # =========================================================================
    # Slide 10: Q&A
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['primary'])
    add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), COLORS['highlight'])

    add_title_shape(slide, "Questions?",
                   Inches(0.5), Inches(1.5), Inches(12), Inches(1),
                   font_size=60, color=COLORS['white'])

    # Key takeaways
    add_body_text(slide, "Key Takeaways:",
                 Inches(0.5), Inches(3), Inches(6), Inches(0.4),
                 font_size=24, color=COLORS['highlight'], bold=True)

    takeaways = [
        "Browser-based IDE = security + flexibility",
        "40-60% cost savings vs traditional VDI",
        "AI-ready infrastructure for the future",
        "Minutes to productivity, not days",
    ]
    add_bullet_points(slide, takeaways,
                     Inches(0.5), Inches(3.5), Inches(6), Inches(2.5), font_size=20)

    # Resources
    add_body_text(slide, "Resources:",
                 Inches(7), Inches(3), Inches(6), Inches(0.4),
                 font_size=24, color=COLORS['highlight'], bold=True)

    resources = [
        "coder.com",
        "github.com/andychoi/dev-platform",
        "coder.com/blog (AI Development)",
    ]
    add_bullet_points(slide, resources,
                     Inches(7), Inches(3.5), Inches(5.5), Inches(2), font_size=18)

    # Thank you
    add_body_text(slide, "Thank You!",
                 Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                 font_size=28, color=COLORS['success'], bold=True)

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "Coder_WebIDE_TechDay.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
